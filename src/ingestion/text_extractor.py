"""
Text extraction module for various file formats.

Supports: PDF, DOCX, DOC, XLSX, XLS, PPTX, HWP, HWPX, CELL, ZIP, HTML, TXT, CSV.
Uses a dispatch pattern to route extraction by file extension.
"""

from __future__ import annotations

import logging
import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Callable

from src.config import settings

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Type alias for extractor functions
# ---------------------------------------------------------------------------
ExtractorFn = Callable[[Path], str]


# ---------------------------------------------------------------------------
# Individual format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(file_path: Path) -> str:
    """Extract text from a PDF using PyMuPDF (fitz)."""
    import fitz  # PyMuPDF

    texts: list[str] = []
    with fitz.open(str(file_path)) as doc:
        for page_num, page in enumerate(doc, start=1):
            page_text = page.get_text()
            if page_text and page_text.strip():
                texts.append(page_text.strip())
    return "\n\n".join(texts)


def _extract_docx(file_path: Path) -> str:
    """Extract text from a DOCX file, including paragraphs and table cells."""
    from docx import Document

    doc = Document(str(file_path))
    parts: list[str] = []

    # Paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Tables
    for table_idx, table in enumerate(doc.tables, start=1):
        table_rows: list[str] = []
        for row in table.rows:
            cell_texts = [cell.text.strip() for cell in row.cells]
            table_rows.append("\t".join(cell_texts))
        if table_rows:
            parts.append(f"[Table {table_idx}]\n" + "\n".join(table_rows))

    return "\n\n".join(parts)


def _extract_xlsx(file_path: Path) -> str:
    """Extract text from an XLSX file, iterating sheets -> rows -> cells."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(v) if v is not None else "" for v in row]
            # Skip completely empty rows
            if any(v.strip() for v in cell_values):
                rows.append("\t".join(cell_values))
        if rows:
            parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(parts)


def _extract_xls(file_path: Path) -> str:
    """Extract text from a legacy .xls file using xlrd."""
    import xlrd

    wb = xlrd.open_workbook(str(file_path))
    parts: list[str] = []

    for sheet in wb.sheets():
        rows: list[str] = []
        for row_idx in range(sheet.nrows):
            cell_values = [
                str(sheet.cell_value(row_idx, col_idx))
                for col_idx in range(sheet.ncols)
            ]
            if any(v.strip() for v in cell_values):
                rows.append("\t".join(cell_values))
        if rows:
            parts.append(f"Sheet: {sheet.name}\n" + "\n".join(rows))

    return "\n\n".join(parts)


def _extract_doc(file_path: Path) -> str:
    """Extract text from a legacy .doc file.

    Strategy:
      1. Try textutil (macOS built-in).
      2. Try antiword (Linux).
      3. Fall back to raw binary text extraction.
    """
    # Attempt 1: textutil (macOS)
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(file_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except FileNotFoundError:
        pass
    except subprocess.TimeoutExpired:
        logger.warning("textutil timed out for %s", file_path)

    # Attempt 2: antiword (Linux)
    try:
        result = subprocess.run(
            ["antiword", str(file_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except FileNotFoundError:
        pass
    except subprocess.TimeoutExpired:
        logger.warning("antiword timed out for %s", file_path)

    # Attempt 3: raw binary text extraction
    try:
        raw = file_path.read_bytes()
        # .doc files often contain readable text mixed with binary.
        # Extract runs of printable text (Korean + ASCII).
        text = raw.decode("utf-8", errors="ignore")
        # Filter to lines with at least some Korean or meaningful ASCII
        lines: list[str] = []
        for line in text.splitlines():
            cleaned = line.strip()
            if len(cleaned) >= 4 and any(
                "\uAC00" <= ch <= "\uD7A3" or ch.isascii() and ch.isalpha()
                for ch in cleaned
            ):
                lines.append(cleaned)
        if lines:
            return "\n".join(lines)
    except Exception:
        logger.debug("Raw text extraction failed for %s", file_path)

    logger.error("All DOC extraction methods failed for %s", file_path)
    return ""


def _extract_pptx(file_path: Path) -> str:
    """Extract text from a PPTX file, iterating slides -> shapes."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"Slide {slide_idx}:\n" + "\n".join(slide_texts))

    return "\n\n".join(parts)


def _extract_hwp(file_path: Path) -> str:
    """
    Extract text from an HWP (Hangul Word Processor) file.

    Strategy:
      1. Try gethwp.read_hwp() first.
      2. Fall back to hwp5txt CLI.
    """
    # Attempt 1: gethwp.read_hwp
    try:
        from gethwp import read_hwp

        text = read_hwp(str(file_path))
        if text and text.strip():
            # Clean up common OLE artifacts (binary garbage chars)
            cleaned = "".join(
                ch for ch in text
                if ch.isprintable() or ch in ("\n", "\t", "\r")
            )
            if cleaned.strip():
                return cleaned.strip()
    except Exception as exc:
        logger.debug("gethwp.read_hwp failed for %s: %s. Trying hwp5txt CLI.", file_path, exc)

    # Attempt 2: hwp5txt CLI
    try:
        result = subprocess.run(
            ["hwp5txt", str(file_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
        if result.stderr:
            logger.warning("hwp5txt stderr for %s: %s", file_path, result.stderr[:500])
    except FileNotFoundError:
        logger.warning("hwp5txt not found in PATH. Cannot extract HWP: %s", file_path)
    except subprocess.TimeoutExpired:
        logger.warning("hwp5txt timed out for %s", file_path)

    logger.error("All HWP extraction methods failed for %s", file_path)
    return ""


def _extract_hwpx(file_path: Path) -> str:
    """
    Extract text from an HWPX file.

    Strategy:
      1. Try gethwp library first.
      2. Fall back to ZIP/XML parsing (HWPX is a ZIP containing XML with hp:t text tags).
    """
    # Attempt 1: gethwp.read_hwpx
    try:
        from gethwp import read_hwpx

        text = read_hwpx(str(file_path))
        if text and text.strip():
            return text.strip()
    except Exception as exc:
        logger.debug("gethwp.read_hwpx failed for %s: %s. Trying ZIP/XML parsing.", file_path, exc)

    # Attempt 2: ZIP/XML parsing
    try:
        import xml.etree.ElementTree as ET

        texts: list[str] = []
        with zipfile.ZipFile(str(file_path), "r") as zf:
            # HWPX stores content in Contents/section*.xml files
            xml_entries = sorted(
                name for name in zf.namelist()
                if name.startswith("Contents/sec") and name.endswith(".xml")
            )
            if not xml_entries:
                # Some variants may use different paths
                xml_entries = sorted(
                    name for name in zf.namelist()
                    if name.endswith(".xml") and "section" in name.lower()
                )

            for entry_name in xml_entries:
                xml_data = zf.read(entry_name)
                root = ET.fromstring(xml_data)
                # Find all hp:t elements (text runs) across all namespaces
                for elem in root.iter():
                    tag = elem.tag
                    # Match tags ending with }t or exactly 't' under hp namespace
                    if tag.endswith("}t") or tag == "t":
                        if elem.text and elem.text.strip():
                            texts.append(elem.text.strip())

        if texts:
            return "\n".join(texts)
    except zipfile.BadZipFile:
        logger.warning("HWPX file is not a valid ZIP: %s", file_path)
    except Exception as exc:
        logger.warning("HWPX ZIP/XML parsing failed for %s: %s", file_path, exc)

    logger.error("All HWPX extraction methods failed for %s", file_path)
    return ""


def _extract_cell(file_path: Path) -> str:
    """
    Extract text from a Hancom CELL file by copying to a temp .xlsx and
    using openpyxl. CELL format is largely xlsx-compatible.
    """
    tmp_dir = None
    try:
        tmp_dir = tempfile.mkdtemp(prefix="cell_extract_")
        tmp_xlsx = Path(tmp_dir) / (file_path.stem + ".xlsx")
        shutil.copy2(str(file_path), str(tmp_xlsx))
        return _extract_xlsx(tmp_xlsx)
    except Exception as exc:
        logger.warning("CELL extraction (xlsx fallback) failed for %s: %s", file_path, exc)
        return ""
    finally:
        if tmp_dir and os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir, ignore_errors=True)


def _extract_html(file_path: Path) -> str:
    """Extract text from an HTML file, removing scripts, styles, and signatures."""
    from bs4 import BeautifulSoup

    raw = _read_with_encoding_fallback(file_path)
    if not raw:
        return ""

    soup = BeautifulSoup(raw, "html.parser")

    # Remove unwanted elements
    for tag_name in ("script", "style", "noscript", "header", "footer", "nav"):
        for tag in soup.find_all(tag_name):
            tag.decompose()

    # Remove common email signature patterns
    for sig_class in ("signature", "gmail_signature", "email-signature", "mail_signature"):
        for tag in soup.find_all(class_=sig_class):
            tag.decompose()
        for tag in soup.find_all(id=sig_class):
            tag.decompose()

    # Remove divs that look like disclaimers
    for div in soup.find_all("div"):
        div_text = div.get_text(strip=True).lower()
        if any(
            keyword in div_text
            for keyword in ("면책", "disclaimer", "confidential", "본 메일은")
        ):
            if len(div_text) < 500:
                div.decompose()

    text = soup.get_text(separator="\n")
    # Collapse multiple blank lines
    lines = [line.strip() for line in text.splitlines()]
    cleaned = "\n".join(line for line in lines if line)
    return cleaned


def _extract_plain_text(file_path: Path) -> str:
    """Read plain text / CSV files with encoding fallback."""
    return _read_with_encoding_fallback(file_path)


# ---------------------------------------------------------------------------
# Archive extraction
# ---------------------------------------------------------------------------

def extract_files_from_archive(
    file_path: Path,
    depth: int = 0,
) -> list[tuple[str, str]]:
    """
    Extract text from all supported files inside a ZIP archive.

    Args:
        file_path: Path to the ZIP file.
        depth: Current recursion depth (to prevent zip bombs).

    Returns:
        List of (internal_path, extracted_text) tuples.
    """
    max_depth = settings.max_zip_depth
    max_size = settings.max_zip_extracted_size_mb * 1024 * 1024  # bytes

    if depth > max_depth:
        logger.warning(
            "Max ZIP depth (%d) exceeded for %s at depth %d. Skipping.",
            max_depth, file_path, depth,
        )
        return []

    results: list[tuple[str, str]] = []
    tmp_dir = None

    try:
        with zipfile.ZipFile(str(file_path), "r") as zf:
            # Check for encryption
            for info in zf.infolist():
                if info.flag_bits & 0x1:
                    logger.warning("Encrypted ZIP detected, skipping: %s", file_path)
                    return []

            # Check total extracted size
            total_size = sum(info.file_size for info in zf.infolist())
            if total_size > max_size:
                logger.warning(
                    "ZIP extracted size (%d bytes) exceeds limit (%d bytes). Skipping: %s",
                    total_size, max_size, file_path,
                )
                return []

            tmp_dir = tempfile.mkdtemp(prefix="zip_extract_")
            zf.extractall(tmp_dir)

        # Process extracted files
        supported = set(settings.supported_extensions)
        tmp_path = Path(tmp_dir)

        for extracted_file in sorted(tmp_path.rglob("*")):
            if not extracted_file.is_file():
                continue

            ext = extracted_file.suffix.lower()
            if ext not in supported:
                continue

            # Compute the internal path relative to the temp directory
            internal_path = str(extracted_file.relative_to(tmp_path))

            if ext == ".zip":
                # Recurse into nested ZIPs
                nested_results = extract_files_from_archive(
                    extracted_file, depth=depth + 1
                )
                for nested_path, nested_text in nested_results:
                    combined_path = f"{internal_path}/{nested_path}"
                    results.append((combined_path, nested_text))
            else:
                text = extract_text(extracted_file, file_extension=ext)
                if text.strip():
                    results.append((internal_path, text))

    except zipfile.BadZipFile:
        logger.error("Bad ZIP file: %s", file_path)
    except Exception as exc:
        logger.error("ZIP extraction failed for %s: %s", file_path, exc)
    finally:
        if tmp_dir and os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir, ignore_errors=True)

    return results


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

_EXTRACTOR_MAP: dict[str, ExtractorFn] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xls,
    ".pptx": _extract_pptx,
    ".hwp": _extract_hwp,
    ".hwpx": _extract_hwpx,
    ".cell": _extract_cell,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".txt": _extract_plain_text,
    ".csv": _extract_plain_text,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_text(file_path: Path, file_extension: str | None = None) -> str:
    """
    Extract text from a file using the appropriate extractor.

    Args:
        file_path: Path to the file to extract text from.
        file_extension: Optional override for the file extension (including the dot).
                        If not provided, it is inferred from the file path.

    Returns:
        Extracted text as a string, or "" if extraction fails.
    """
    ext = (file_extension or file_path.suffix).lower()

    if ext == ".zip":
        # For ZIP files, concatenate all extracted texts
        results = extract_files_from_archive(file_path)
        if not results:
            return ""
        parts = []
        for internal_path, text in results:
            parts.append(f"[{internal_path}]\n{text}")
        return "\n\n".join(parts)

    extractor = _EXTRACTOR_MAP.get(ext)
    if extractor is None:
        logger.warning("Unsupported file extension '%s' for file: %s", ext, file_path)
        return ""

    try:
        text = extractor(file_path)
        return text if text else ""
    except Exception as exc:
        logger.error(
            "Text extraction failed for %s (ext=%s): %s",
            file_path, ext, exc,
            exc_info=True,
        )
        return ""


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def _read_with_encoding_fallback(file_path: Path) -> str:
    """
    Read a text file trying multiple encodings in order.

    Tries: UTF-8 -> EUC-KR -> CP949 -> Latin-1 (latin-1 never fails).
    """
    encodings = ["utf-8", "euc-kr", "cp949", "latin-1"]
    for encoding in encodings:
        try:
            return file_path.read_text(encoding=encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
        except Exception as exc:
            logger.error("Failed to read %s with encoding %s: %s", file_path, encoding, exc)
            return ""
    # Should never reach here since latin-1 accepts all byte values
    return ""
